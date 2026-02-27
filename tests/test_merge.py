from pathlib import Path

from PIL import Image
from pptx import Presentation

from app.ppt_merge import merge_with_template


def _build_template(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "模板首页"
    prs.save(str(path))


def _build_source(path: Path) -> None:
    img_path = path.parent / "tmp.png"
    Image.new("RGB", (200, 120), (43, 108, 176)).save(str(img_path))

    prs = Presentation()
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    slide1.shapes.add_textbox(500000, 400000, 6000000, 800000).text = "源文件内容 A"
    slide1.shapes.add_picture(str(img_path), 500000, 1500000, width=3000000)

    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    slide2.shapes.add_textbox(400000, 400000, 5000000, 800000).text = "源文件内容 B"
    prs.save(str(path))


def test_merge_with_template(tmp_path: Path):
    template = tmp_path / "template.pptx"
    source = tmp_path / "source.pptx"
    output = tmp_path / "output.pptx"

    _build_template(template)
    _build_source(source)

    report = merge_with_template(template, [source], output)

    merged = Presentation(str(output))
    assert output.exists()
    assert len(merged.slides) == 3
    assert report.total_source_files == 1
    assert report.total_source_slides == 2
    assert report.imported_slides == 2
