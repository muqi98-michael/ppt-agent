from pathlib import Path

from fastapi.testclient import TestClient
from pptx import Presentation

from app.main import app


def _create_ppt(path: Path, title: str) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(300000, 300000, 4000000, 600000).text = title
    prs.save(str(path))


def test_merge_endpoint(tmp_path: Path):
    template_path = tmp_path / "template.pptx"
    source1_path = tmp_path / "source1.pptx"
    source2_path = tmp_path / "source2.pptx"
    _create_ppt(template_path, "模板")
    _create_ppt(source1_path, "内容1")
    _create_ppt(source2_path, "内容2")

    client = TestClient(app)
    with template_path.open("rb") as template_f, source1_path.open("rb") as source1_f, source2_path.open("rb") as source2_f:
        response = client.post(
            "/merge",
            files=[
                ("template", ("template.pptx", template_f, "application/vnd.openxmlformats-officedocument.presentationml.presentation")),
                ("sources", ("source1.pptx", source1_f, "application/vnd.openxmlformats-officedocument.presentationml.presentation")),
                ("sources", ("source2.pptx", source2_f, "application/vnd.openxmlformats-officedocument.presentationml.presentation")),
            ],
        )

    assert response.status_code == 200
    assert response.headers["content-type"] == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    assert int(response.headers["x-merge-imported-slides"]) == 2
    assert response.content[:2] == b"PK"
