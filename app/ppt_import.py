"""PPT 自动入库：按页提取文本、按章节拆分、生成章节 PPT"""
from __future__ import annotations

import base64
import re
import zipfile
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from tempfile import TemporaryDirectory

from pptx import Presentation

from app.ppt_merge import _copy_slide_content, _pick_blank_layout
from app.research import summarize_chapter_contents


@dataclass
class SlideText:
    """单页幻灯片的文本内容"""
    index: int  # 0-based
    title: str
    body: str
    full_text: str


@dataclass
class Chapter:
    """章节信息"""
    title: str
    slide_indices: list[int]  # 0-based
    content: str  # 合并的文本内容
    summary: str  # 大模型提炼摘要
    ppt_base64: str  # 该章节 PPT 的 base64


def extract_slide_text(slide, index: int) -> SlideText:
    """从单页幻灯片提取文本"""
    title_parts = []
    body_parts = []

    for shape in slide.shapes:
        if not hasattr(shape, "text") or not shape.text:
            continue
        text = shape.text.strip()
        if not text:
            continue

        # 简单启发式：标题通常在顶部，文本较短；正文在下方
        is_placeholder = getattr(shape, "is_placeholder", False)
        ph_type = getattr(shape.placeholder_format, "type", None) if is_placeholder else None

        if is_placeholder and ph_type in (1, 2, 18):  # title, centerTitle, etc
            title_parts.append(text)
        elif not title_parts and len(text) < 80:
            title_parts.append(text)
        else:
            body_parts.append(text)

    title = "\n".join(title_parts) if title_parts else ""
    body = "\n".join(body_parts) if body_parts else ""
    full = f"{title}\n{body}".strip() if title or body else ""

    return SlideText(index=index, title=title, body=body, full_text=full)


def extract_all_text(pres: Presentation) -> list[SlideText]:
    """从演示文稿所有页面提取文本"""
    return [extract_slide_text(slide, i) for i, slide in enumerate(pres.slides)]


def detect_chapter_boundaries(slides_text: list[SlideText]) -> list[tuple[int, int, str]]:
    """检测章节边界，返回 [(start_idx, end_idx, chapter_title), ...]
    章节起始页（如 01、02 大数字+标题）及其后续内容为一章，拆成独立 PPT。
    """
    if not slides_text:
        return []

    # 章节起始页模式：
    # 1) 01、02、03 等大数字（单独或后跟标题）
    # 2) 第X章、第X节、第一章、Part 1、Chapter 1、一、二、三 等
    section_number_pattern = re.compile(r"^0?\d{1,2}\b")  # 01、02、1、2
    chapter_pattern = re.compile(
        r"^(第[一二三四五六七八九十百千\d]+[章节部分]|"
        r"第\d+[章节部分]|"
        r"Part\s*\d+|"
        r"Chapter\s*\d+|"
        r"第\s*\d+\s*[章节]|"
        r"^[一二三四五六七八九十]+[、．.]|"
        r"^\d+[\.、．]\s*\S)",
        re.IGNORECASE,
    )

    def _clean_title(s: str) -> str:
        if not s:
            return ""
        s = re.sub(r"[\s\x00-\x1f]+", " ", s).strip()[:80]
        return s

    # 常见章节标题结尾（方案介绍、发展趋势、整体方案等）
    section_title_suffixes = ("介绍", "方案", "趋势", "发展", "概述", "总览", "布局", "蓝图")

    def _is_chapter_start(st: SlideText) -> bool:
        """判断是否为章节起始页（该页+后续为一章）
        如：01/02 大数字+标题、第X章、或仅含简短标题的章节分隔页
        """
        text = (st.title or "") + "\n" + (st.body or "")
        lines = [ln.strip() for ln in text.split("\n") if ln.strip()]
        first_line = lines[0] if lines else ""
        full_len = len(st.full_text)
        # 1) 01、02 等大数字开头（含整页仅数字的情况）
        if section_number_pattern.match(first_line) and full_len < 500:
            return True
        # 2) 任意行以 01、02 开头（大数字可能在第二行）
        for ln in lines[:3]:
            if section_number_pattern.match(ln) and full_len < 500:
                return True
        # 3) 传统章节模式
        if chapter_pattern.match(st.title.strip() or "") and full_len < 500:
            return True
        # 4) 章节分隔页：仅含简短标题，无长正文
        if st.title and 4 <= len(st.title) <= 80:
            if not st.body or len(st.body) < 80:
                if full_len < 250:
                    return True
        # 5) 标题以常见章节词结尾（如「xxx整体方案介绍」），且内容较短
        title_clean = (st.title or "").strip()
        if title_clean and any(title_clean.endswith(s) for s in section_title_suffixes):
            if full_len < 300 and (not st.body or len(st.body) < 100):
                return True
        # 6) 内容极短（8-35字）且含中文，多为章节分隔页（如仅「02」+标题）
        if 8 <= full_len <= 35 and re.search(r"[\u4e00-\u9fff]{4,}", st.full_text):
            return True
        return False

    def _extract_chapter_title(st: SlideText) -> str:
        """从章节起始页提取标题（去掉纯数字行，取有意义的标题）"""
        lines = (st.title or "").split("\n") + (st.body or "").split("\n")
        for line in lines:
            line = line.strip()
            if not line:
                continue
            if section_number_pattern.match(line) and len(line) <= 3:
                continue  # 跳过纯 "01"、"02"
            if len(line) > 2 and not line.isdigit():
                return _clean_title(line)
        return _clean_title(st.title or st.full_text) or "未命名章节"

    chapters: list[tuple[int, int, str]] = []
    current_start = 0
    current_title = _extract_chapter_title(slides_text[0])

    for i, st in enumerate(slides_text):
        if _is_chapter_start(st):
            if i > current_start:
                chapters.append((current_start, i, current_title))
            current_start = i
            current_title = _extract_chapter_title(st)

    chapters.append((current_start, len(slides_text), current_title))
    return chapters


def create_chapter_ppt(
    source_path: Path,
    slide_indices: list[int],
    output_path: Path,
) -> None:
    """从源 PPT 中提取指定页，生成新 PPT"""
    source = Presentation(str(source_path))
    target_width = int(source.slide_width)
    target_height = int(source.slide_height)

    dest = Presentation()
    dest.slide_width = target_width
    dest.slide_height = target_height
    blank_layout = _pick_blank_layout(dest)

    for idx in slide_indices:
        if idx >= len(source.slides):
            continue
        src_slide = source.slides[idx]
        dest_slide = dest.slides.add_slide(blank_layout)
        _copy_slide_content(src_slide, dest_slide)

    dest.save(str(output_path))


def process_ppt_import(source_path: Path) -> tuple[list[Chapter], bytes]:
    """处理 PPT 入库：提取文本、拆分章节、生成章节 PPT，并打包为 zip。
    返回 (chapters, zip_bytes)。
    """
    pres = Presentation(str(source_path))
    slides_text = extract_all_text(pres)
    if not slides_text:
        return [], b""

    boundaries = detect_chapter_boundaries(slides_text)
    chapter_drafts: list[dict[str, str | list[int]]] = []
    zip_buffer = BytesIO()

    for i, (start, end, title) in enumerate(boundaries):
        indices = list(range(start, end))
        content_parts = [slides_text[j].full_text for j in indices if slides_text[j].full_text]
        content = "\n\n---\n\n".join(content_parts)
        chapter_drafts.append(
            {
                "title": title,
                "indices": indices,
                "content": content,
                "safe_title": re.sub(r"[^\w\s\u4e00-\u9fff-]", "_", title)[:30] or f"chapter_{i + 1}",
            }
        )

    summaries = summarize_chapter_contents(
        [{"title": str(item["title"]), "content": str(item["content"])} for item in chapter_drafts]
    )
    chapters: list[Chapter] = []

    with TemporaryDirectory() as tmp_dir:
        tmp_path = Path(tmp_dir)
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
            for i, item in enumerate(chapter_drafts):
                title = str(item["title"])
                indices = list(item["indices"])
                content = str(item["content"])
                safe_title = str(item["safe_title"])
                chapter_base_name = f"章节{i + 1}_{safe_title}"
                chapter_pptx = tmp_path / f"{chapter_base_name}.pptx"
                create_chapter_ppt(source_path, indices, chapter_pptx)

                ppt_bytes = chapter_pptx.read_bytes()
                ppt_b64 = base64.b64encode(ppt_bytes).decode("utf-8")
                chapter_summary = summaries[i] if i < len(summaries) else ""
                chapter_md = _build_chapter_markdown(
                    title=title,
                    summary=chapter_summary,
                    content=content,
                    slide_count=len(indices),
                )

                zf.writestr(f"{chapter_base_name}.pptx", ppt_bytes)
                zf.writestr(f"{chapter_base_name}.md", chapter_md.encode("utf-8"))

                chapters.append(
                    Chapter(
                        title=title,
                        slide_indices=indices,
                        content=content,
                        summary=chapter_summary,
                        ppt_base64=ppt_b64,
                    )
                )

    return chapters, zip_buffer.getvalue()


def _build_chapter_markdown(*, title: str, summary: str, content: str, slide_count: int) -> str:
    clean_title = title.strip() or "未命名章节"
    clean_summary = summary.strip() or "暂无摘要"
    clean_content = content.strip() or "暂无正文内容"
    return (
        f"# {clean_title}\n\n"
        f"- 页数: {slide_count}\n\n"
        f"## 概要总结\n\n"
        f"{clean_summary}\n\n"
        f"## 章节正文\n\n"
        f"{clean_content}\n"
    )
