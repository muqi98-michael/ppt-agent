from __future__ import annotations

from copy import deepcopy
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.oxml.xmlchemy import OxmlElement

REL_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"


@dataclass
class MergeReport:
    total_source_files: int
    total_source_slides: int
    imported_slides: int
    layout_adjusted_slides: int


def merge_with_template(
    template_path: str | Path,
    source_paths: list[str | Path],
    output_path: str | Path,
) -> MergeReport:
    """Merge all source slides into template deck and keep source formatting.

    Strategy:
    - Keep the template's existing slides untouched.
    - Append each source slide as a new blank slide.
    - Copy shape XML + relationship references to preserve visual styling.
    - Normalize positions when source slide size differs from template size.
    """
    template = Presentation(str(template_path))
    target_width = int(template.slide_width)
    target_height = int(template.slide_height)
    blank_layout = _pick_blank_layout(template)

    report = MergeReport(
        total_source_files=len(source_paths),
        total_source_slides=0,
        imported_slides=0,
        layout_adjusted_slides=0,
    )

    for source_path in source_paths:
        source = Presentation(str(source_path))
        report.total_source_slides += len(source.slides)

        for src_slide in source.slides:
            dest_slide = template.slides.add_slide(blank_layout)
            _clear_slide_shapes(dest_slide)
            _disable_master_overlay(dest_slide)
            _copy_slide_content(src_slide, dest_slide)

            if _normalize_slide_layout(
                dest_slide,
                src_width=int(source.slide_width),
                src_height=int(source.slide_height),
                dst_width=target_width,
                dst_height=target_height,
            ):
                report.layout_adjusted_slides += 1

            report.imported_slides += 1

    template.save(str(output_path))
    return report


def _pick_blank_layout(pres: Presentation):
    best_layout = pres.slide_layouts[0]
    best_score = float("inf")

    for layout in pres.slide_layouts:
        name = (layout.name or "").lower()
        placeholder_count = sum(1 for shape in layout.shapes if shape.is_placeholder)
        text_count = sum(
            1
            for shape in layout.shapes
            if hasattr(shape, "text") and shape.text and shape.text.strip()
        )
        shape_count = len(layout.shapes)

        name_bonus = 0
        if "blank" in name or "空白" in name:
            name_bonus -= 200
        if "标题幻灯片" in name:
            name_bonus -= 20

        score = placeholder_count * 100 + text_count * 10 + shape_count + name_bonus
        if score < best_score:
            best_score = score
            best_layout = layout

    return best_layout


def _copy_slide_content(src_slide, dest_slide) -> None:
    sp_tree = dest_slide.shapes._spTree
    rel_map: dict[tuple[str, str], str] = {}

    # Copy source master/layout visual artifacts first for pages that rely on
    # non-placeholder layout design (e.g. chapter cover background).
    _copy_shapes_from_container(
        source_shapes=src_slide.slide_layout.shapes,
        source_part=src_slide.slide_layout.part,
        dest_slide=dest_slide,
        sp_tree=sp_tree,
        rel_map=rel_map,
        strip_placeholders=False,
        skip_placeholders=True,
    )

    # Copy slide-level shapes last so content stays on top.
    _copy_shapes_from_container(
        source_shapes=src_slide.shapes,
        source_part=src_slide.part,
        dest_slide=dest_slide,
        sp_tree=sp_tree,
        rel_map=rel_map,
        strip_placeholders=True,
        skip_placeholders=False,
    )

    # Copy effective source background (slide -> layout -> master).
    src_bg = src_slide.element.cSld.bg
    if src_bg is None:
        src_bg = src_slide.slide_layout.element.cSld.bg
    if src_bg is None:
        src_bg = src_slide.slide_layout.slide_master.element.cSld.bg
    if src_bg is not None:
        dest_cSld = dest_slide.element.cSld
        if dest_cSld.bg is not None:
            dest_cSld.remove(dest_cSld.bg)
        dest_cSld.insert(0, deepcopy(src_bg))


def _copy_shapes_from_container(
    source_shapes,
    source_part,
    dest_slide,
    sp_tree,
    rel_map: dict[tuple[str, str], str],
    *,
    strip_placeholders: bool,
    skip_placeholders: bool,
) -> None:
    for shape in source_shapes:
        if skip_placeholders and getattr(shape, "is_placeholder", False):
            continue

        copied = deepcopy(shape.element)
        if strip_placeholders:
            placeholder_type, placeholder_idx = _get_placeholder_info(copied)
            if placeholder_type in {"title", "ctrTitle"}:
                _materialize_placeholder_geometry(copied, source_part, placeholder_type, placeholder_idx)
                title_size = _resolve_effective_title_size(shape, source_part, placeholder_type, placeholder_idx)
            else:
                title_size = None
            placeholder_type = _strip_placeholder_binding(copied)
            if placeholder_type in {"title", "ctrTitle"}:
                _apply_fallback_title_style(copied, title_size)

        _remap_relationship_ids(copied, source_part, dest_slide, rel_map)
        sp_tree.insert_element_before(copied, "p:extLst")


def _apply_fallback_title_style(shape_el, fallback_size: str | None = None) -> None:
    """Ensure title text remains visible after detaching placeholder binding."""
    target_size = fallback_size or "2400"
    for node in shape_el.iter():
        if not (node.tag.endswith("}rPr") or node.tag.endswith("}endParaRPr")):
            continue
        if node.get("sz") is None:
            node.set("sz", target_size)
        if not any(child.tag.endswith("}solidFill") for child in node):
            solid_fill = OxmlElement("a:solidFill")
            scheme = OxmlElement("a:schemeClr")
            scheme.set("val", "tx1")
            solid_fill.append(scheme)
            node.append(solid_fill)


def _get_placeholder_info(shape_el) -> tuple[str | None, str | None]:
    for node in shape_el.iter():
        if node.tag.endswith("}ph"):
            return node.get("type"), node.get("idx")
    return None, None


def _materialize_placeholder_geometry(shape_el, source_part, ph_type: str | None, ph_idx: str | None) -> None:
    """Copy xfrm from source layout/master placeholder when slide placeholder has no own geometry."""
    sp_pr = shape_el.find(".//p:spPr", shape_el.nsmap)
    if sp_pr is None:
        return
    if sp_pr.find(".//a:xfrm", shape_el.nsmap) is not None:
        return

    source_slide = getattr(source_part, "slide", None)
    if source_slide is None:
        return

    placeholder_source = _find_placeholder_shape(source_slide.slide_layout, ph_type, ph_idx)
    if placeholder_source is None:
        placeholder_source = _find_placeholder_shape(source_slide.slide_layout.slide_master, ph_type, ph_idx)
    if placeholder_source is None:
        return

    src_sp_pr = placeholder_source.element.find(".//p:spPr", shape_el.nsmap)
    if src_sp_pr is None:
        return
    src_xfrm = src_sp_pr.find(".//a:xfrm", shape_el.nsmap)
    if src_xfrm is None:
        return
    sp_pr.append(deepcopy(src_xfrm))


def _find_placeholder_shape(container, ph_type: str | None, ph_idx: str | None):
    for shape in container.shapes:
        if not getattr(shape, "is_placeholder", False):
            continue
        c_type, c_idx = _get_placeholder_info(shape.element)
        if c_type != ph_type:
            continue
        if (c_idx or "0") != (ph_idx or "0"):
            continue
        return shape
    return None


def _resolve_effective_title_size(
    source_shape,
    source_part,
    ph_type: str | None,
    ph_idx: str | None,
) -> str | None:
    """Get title font size from source slide first, then layout/master placeholders."""
    size = _find_first_text_size(source_shape.element)
    if size is not None:
        return size

    source_slide = getattr(source_part, "slide", None)
    if source_slide is None:
        return None

    layout_ph = _find_placeholder_shape(source_slide.slide_layout, ph_type, ph_idx)
    if layout_ph is not None:
        size = _find_first_text_size(layout_ph.element)
        if size is not None:
            return size

    master_ph = _find_placeholder_shape(source_slide.slide_layout.slide_master, ph_type, ph_idx)
    if master_ph is not None:
        size = _find_first_text_size(master_ph.element)
        if size is not None:
            return size

    return None


def _find_first_text_size(shape_el) -> str | None:
    for node in shape_el.iter():
        if node.tag.endswith("}rPr") and node.get("sz"):
            return node.get("sz")
    for node in shape_el.iter():
        if node.tag.endswith("}endParaRPr") and node.get("sz"):
            return node.get("sz")
    return None


def _clear_slide_shapes(slide) -> None:
    """Remove all layout-provided shapes from a newly created slide."""
    sp_tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        sp_tree.remove(shape.element)


def _disable_master_overlay(slide) -> None:
    """Hide template master/layout artifacts on merged pages."""
    slide.element.set("showMasterSp", "0")
    slide.element.set("showMasterPhAnim", "0")
    try:
        slide.follow_master_background = False
    except Exception:
        # Some decks may not expose this setter consistently.
        pass


def _strip_placeholder_binding(shape_el) -> str | None:
    """Remove placeholder marker so copied shape won't inherit template placeholder style."""
    placeholder_type = None
    for node in list(shape_el.iter()):
        if node.tag.endswith("}ph"):
            placeholder_type = node.get("type")
            parent = node.getparent()
            if parent is not None:
                parent.remove(node)
    return placeholder_type


def _remap_relationship_ids(shape_el, source_part, dest_slide, rel_map: dict[tuple[str, str], str]) -> None:
    nodes_to_remove = []

    for node in shape_el.iter():
        for attr_name, attr_value in list(node.attrib.items()):
            if not attr_name.startswith(REL_NS):
                continue

            old_rid = attr_value
            rel_key = (str(source_part.partname), old_rid)
            if rel_key in rel_map:
                node.set(attr_name, rel_map[rel_key])
                continue

            rel = source_part.rels.get(old_rid)
            if rel is None:
                continue

            # Tags are PowerPoint metadata artifacts. Dropping them avoids
            # duplicate package-part conflicts while keeping visual content intact.
            if rel.reltype.endswith("/tags"):
                nodes_to_remove.append(node)
                break

            if rel.reltype.endswith("/image"):
                try:
                    image_blob = rel.target_part.blob
                    _, new_rid = dest_slide.part.get_or_add_image_part(BytesIO(image_blob))
                    rel_map[rel_key] = new_rid
                    node.set(attr_name, new_rid)
                    continue
                except Exception:
                    # For uncommon/legacy image encodings fallback to raw rel copy.
                    pass

            if rel.is_external:
                new_rid = dest_slide.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
            else:
                new_rid = dest_slide.part.rels.get_or_add(rel.reltype, rel.target_part)

            rel_map[rel_key] = new_rid
            node.set(attr_name, new_rid)

    for node in nodes_to_remove:
        parent = node.getparent()
        if parent is not None:
            parent.remove(node)


def _normalize_slide_layout(
    slide,
    src_width: int,
    src_height: int,
    dst_width: int,
    dst_height: int,
) -> bool:
    if src_width == dst_width and src_height == dst_height:
        return False

    scale_x = dst_width / src_width
    scale_y = dst_height / src_height

    for shape in slide.shapes:
        if hasattr(shape, "left"):
            shape.left = int(shape.left * scale_x)
        if hasattr(shape, "top"):
            shape.top = int(shape.top * scale_y)
        if hasattr(shape, "width"):
            shape.width = max(1, int(shape.width * scale_x))
        if hasattr(shape, "height"):
            shape.height = max(1, int(shape.height * scale_y))

    return True
