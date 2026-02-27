from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

from app.research import ResearchResult, SectionResult


def append_research_slides(
    template_path: str | Path,
    output_path: str | Path,
    research: ResearchResult,
) -> None:
    prs = Presentation(str(template_path))
    layout = _pick_layout(prs)

    _add_research_slide(
        prs,
        layout,
        title=f"{research.industry}：趋势、痛点与IT规划",
        subtitle="基于公开网络信息自动整理",
        sections=research.industry_sections,
    )
    _add_research_slide(
        prs,
        layout,
        title=f"{research.customer}：规划、痛点与IT需求",
        subtitle="基于官网与公开新闻自动整理",
        sections=research.customer_sections,
    )

    prs.save(str(output_path))


def _pick_layout(prs: Presentation):
    """优先使用「白色内页」版式。"""
    for layout in prs.slide_layouts:
        name = (layout.name or "").strip()
        if name == "白色内页":
            return layout
    for layout in prs.slide_layouts:
        name = (layout.name or "").strip()
        if "白色内页" in name:
            return layout
    for layout in prs.slide_layouts:
        if "blank" in (layout.name or "").lower() or "空白" in (layout.name or ""):
            return layout
    return prs.slide_layouts[0]


def _add_research_slide(prs: Presentation, layout, title: str, subtitle: str, sections: list[SectionResult]) -> None:
    slide = prs.slides.add_slide(layout)
    _clear_slide_shapes(slide)

    title_box = slide.shapes.add_textbox(500000, 280000, 11200000, 520000)
    tf = title_box.text_frame
    tf.word_wrap = True
    tr = tf.paragraphs[0].add_run()
    tr.text = title
    tr.font.size = Pt(24)
    tr.font.bold = True

    sub_box = slide.shapes.add_textbox(500000, 820000, 11200000, 240000)
    sf_sub = sub_box.text_frame
    sf_sub.word_wrap = True
    sr = sf_sub.paragraphs[0].add_run()
    sr.text = subtitle
    sr.font.size = Pt(14)

    card_w = 3560000
    lefts = [500000, 4300000, 8100000]
    top = 1200000
    card_h = 4500000

    for idx, section in enumerate(sections[:3]):
        x = lefts[idx]
        card = slide.shapes.add_shape(1, x, top, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = _rgb("F8FAFC")
        card.line.color.rgb = _rgb("E2E8F0")

        header = slide.shapes.add_textbox(x + 120000, top + 80000, card_w - 240000, 320000)
        hf = header.text_frame
        hf.word_wrap = True
        hr = hf.paragraphs[0].add_run()
        hr.text = section.title
        hr.font.size = Pt(16)
        hr.font.bold = True

        bullet_box = slide.shapes.add_textbox(x + 120000, top + 440000, card_w - 240000, 3200000)
        bf = bullet_box.text_frame
        bf.word_wrap = True
        for bi, bullet in enumerate(section.bullets[:5]):
            p = bf.add_paragraph() if bi > 0 else bf.paragraphs[0]
            p.text = f"• {bullet}"
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(6)

        if section.sources:
            src_box = slide.shapes.add_textbox(x + 120000, top + 3760000, card_w - 240000, 680000)
            sframe = src_box.text_frame
            sframe.word_wrap = True
            sp0 = sframe.paragraphs[0]
            r0 = sp0.add_run()
            r0.text = "来源: "
            r0.font.bold = True
            r0.font.size = Pt(10)
            for src in section.sources[:3]:
                sp = sframe.add_paragraph()
                sp.text = src
                sp.font.size = Pt(9)


def _clear_slide_shapes(slide) -> None:
    sp_tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        sp_tree.remove(shape.element)


def _rgb(hex_color: str):
    from pptx.dml.color import RGBColor

    return RGBColor.from_string(hex_color)
